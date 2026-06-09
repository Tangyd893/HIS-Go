"""Generate two HIS-Go demo PPTs: Admin + Patient."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Color palette ──
BLUE = RGBColor(0x18, 0x90, 0xFF)
DARK_BLUE = RGBColor(0x0A, 0x5C, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF0, 0xF5, 0xFF)
GREEN = RGBColor(0x00, 0xA8, 0x70)
RED = RGBColor(0xE8, 0x3E, 0x3E)
ORANGE = RGBColor(0xFA, 0xAD, 0x14)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, bg_color=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, bg_color)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9.3), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(subtitle):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xD6, 0xEA, 0xFF)
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, section_title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9.3), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, note="", demo_table=None):
    """Add a standard content slide with title + bullets + optional demo table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Top bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    # Bullets
    top_y = Inches(1.2)
    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), top_y, Inches(11.5), Inches(4.5) if demo_table else Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(18)
            p2.font.color.rgb = DARK
            p2.space_after = Pt(8)
        top_y = Inches(1.2 + len(bullets) * 0.38 + 0.2)
    # Demo table
    if demo_table:
        headers, rows = demo_table
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_w = Inches(11)
        tbl_h = Inches(0.4 * n_rows)
        left = Inches(1)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top_y, tbl_w, tbl_h)
        table = table_shape.table
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = val
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(13)
                    para.font.color.rgb = DARK
    # Note
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6))
        nf = note_box.text_frame
        nf.word_wrap = True
        np = nf.paragraphs[0]
        np.text = f"💬 {note}"
        np.font.size = Pt(12)
        np.font.italic = True
        np.font.color.rgb = GRAY
    return slide


def generate_admin_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── 1. Cover ──
    add_title_slide(prs,
        "HIS-Go 医院信息系统",
        ["Web 管理端演示", "全链路 Go 微服务重构 · 门诊业务闭环", "GitHub: Tangyd893/HIS-Go"])

    # ── 2. Background ──
    add_content_slide(prs,
        "为什么用 Go 重构 HIS？",
        [
            "原系统：Spring Cloud Alibaba + Vue（Hospital-Information-System）",
            "目标：17+ 微服务迁移至 Go，保持业务逻辑与分库设计",
            "收益：更低资源占用、更简部署、gRPC 强类型通信",
            "范围：院内诊疗 + 院外患者服务",
        ],
        note="强调「重构而非从零」，降低评审对业务完整性的质疑。")

    # ── 3. Architecture ──
    add_content_slide(prs,
        "系统总体架构",
        [
            "🖥️ 客户端：Vue3 管理端（frontend/admin）",
            "🔀 接入层：Nginx :80 → Gateway :8080（JWT / 限流 / 路由）",
            "⚙️ 服务层：18 个 Go 微服务，gRPC 内部通信",
            "💾 数据层：PostgreSQL 分库 · Redis · RabbitMQ",
        ])

    # ── 4. Microservices ──
    add_content_slide(prs,
        "微服务拆分 — Database per Service",
        [
            "🔐 auth 认证 · 👤 user 患者/员工/科室",
            "📅 schedule 排班 · 🏥 registration 挂号",
            "🩺 clinic 门诊诊疗 · 💊 prescription 处方",
            "💰 billing 收费 · 📦 pharmacy 药房 · ⚙️ system 系统",
        ],
        note="每个服务独立数据库，独立 cmd/<name>/main.go 入口。")

    # ── 5. Tech Stack ──
    add_content_slide(prs,
        "技术选型",
        [
            "后端：Go 1.25 · Gin · gRPC · GORM · Wire",
            "前端：Vue 3 · TypeScript · Ant Design Vue · Vite",
            "基础设施：PostgreSQL 17 · Redis 7 · RabbitMQ 4",
            "部署：Docker Compose 多阶段构建",
        ])

    # ── 6. Deployment ──
    add_content_slide(prs,
        "本地演示一键启动",
        [
            "Profile：deploy/compose/demo-admin.yml",
            "含 PG + Redis + RabbitMQ + 10 微服务 + Gateway + Nginx",
            "环境变量：deploy/config/demo.env",
            "管理端静态资源：frontend/admin/dist 由 Nginx 托管",
        ],
        note="docker compose -f deploy/compose/demo-admin.yml --env-file deploy/config/demo.env up -d --build")

    # ── 7. Demo Account ──
    add_content_slide(prs,
        "演示环境说明",
        [
            "入口：http://localhost/admin",
            "管理员账号：demo-admin / demo123",
            "医生账号：demo-doctor / demo123",
            "⚠️ 仅限本地演示，密码不可用于生产",
            "演示数据来自 seed_data.sql（10 患者、5 科室、8 药品）",
        ])

    # ── 8. Business Flow Overview ──
    add_content_slide(prs,
        "门诊主路径 — 六步闭环",
        [
            "① 排班管理 — 发布号源",
            "② 挂号管理 — 预约 + 签到",
            "③ 门诊诊疗 — 主诉/诊断",
            "④ 处方管理 — 开方 + 审核",
            "⑤ 收费结算 — 账单支付",
            "⑥ 药房管理 — 发药出库",
        ])

    # ── 9. Demo 1: Registration ──
    add_content_slide(prs,
        "【现场演示】① 挂号与签到",
        [
            "步骤 1：使用 demo-admin 登录系统",
            "步骤 2：挂号管理 → 查询今天号源 → 挂王小明",
            "步骤 3：挂号记录 → 签到 → 开始接诊",
        ],
        demo_table=(
            ["步骤", "页面", "操作"],
            [
                ["1", "登录", "demo-admin 登录"],
                ["2", "挂号管理", "今天号源 → 挂王小明"],
                ["3", "挂号记录", "签到 → 开始接诊"],
            ]
        ),
        note="患者搜索：姓名模糊匹配；号源「剩余」随挂号扣减。")

    # ── 10. Demo 2: Clinic ──
    add_content_slide(prs,
        "【现场演示】② 诊疗与处方审核",
        [
            "步骤 4：门诊诊疗 → 主诉「发热咳嗽3天」→ 诊断「上呼吸道感染」→ 保存",
            "步骤 5：处方管理 → 开方（阿莫西林胶囊）→ 审核",
        ],
        demo_table=(
            ["步骤", "页面", "操作"],
            [
                ["4", "门诊诊疗", "主诉/诊断 → 保存"],
                ["5", "处方管理", "开方（阿莫西林）→ 审核"],
            ]
        ),
        note="处方状态流转：待审核 → 已审核 → 已发药。")

    # ── 11. Demo 3: Billing + Pharmacy ──
    add_content_slide(prs,
        "【现场演示】③ 收费发药闭环",
        [
            "步骤 6：收费结算 → 待支付 → 收费",
            "步骤 7：药房管理 → 选处方 → 发药",
            "步骤 8：回处方管理 → 确认「已发药」",
        ],
        demo_table=(
            ["步骤", "页面", "操作"],
            [
                ["6", "收费结算", "待支付 → 收费"],
                ["7", "药房管理", "选处方 → 发药"],
                ["8", "处方管理", "确认「已发药」"],
            ]
        ),
        note="支付渠道为演示 Stub；发药扣减库存，闭环完成。")

    # ── 12. Gateway ──
    add_content_slide(prs,
        "API 网关与鉴权",
        [
            "Gateway 路由：/api/auth → auth:8081，/api/user → user:8082 …",
            "JWT 签发：auth 服务；网关校验 Token",
            "健康检查：GET /health",
            "演示阶段 RBAC 简化，生产需完善 RequireRole",
        ])

    # ── 13. Data Layer ──
    add_content_slide(prs,
        "PostgreSQL 分库 + Seed 数据",
        [
            "每服务独立库：his_auth、his_user、his_clinic …",
            "初始化：backend/migrations/ 自动建库建表",
            "演示数据：backend/sql/seed_data.sql（10 患者、5 科室、8 药品）",
            "宿主机 PG 端口：5433（避免与本机冲突）",
        ])

    # ── 14. Frontend ──
    add_content_slide(prs,
        "Vue 3 管理端前端架构",
        [
            "路由：frontend/admin/src/router/index.ts（18 菜单）",
            "状态管理：Pinia store/auth.ts",
            "API 封装：api/client.ts 统一 Bearer Token",
            "演示 Profile 仅 10 个后端 — 未启动菜单显示占位页",
        ],
        note="演示时勿点 EMR/住院/检查等未启动项。")

    # ── 15. Completion Status ──
    add_content_slide(prs,
        "功能完成度（诚实口径）",
        [
            "✅ 可演示：挂号、排班、诊疗、处方、收费、药房、患者/科室",
            "⚠️ 占位 UI：电子病历、住院、检查检验、随访、通知等",
            "🔧 演示级：支付 Stub、计费金额简化、部分详情按钮占位",
        ],
        note="定位为「可脚本化本地演示」，非生产就绪。")

    # ── 16. Quality ──
    add_content_slide(prs,
        "测试与 CI",
        [
            "单元测试：go test ./... 通过",
            "集成测试：testing/api/demo_admin_flow_test.go（需 Docker）",
            "CI：.github/workflows/ci.yml — vet / test / 前端 build / SonarCloud",
            "前端构建：admin + patient 双端 npm run build 校验",
        ])

    # ── 17. Limitations ──
    add_content_slide(prs,
        "后续优化方向",
        [
            "首次部署：Docker 构建耗时、seed 自动化（db-init）",
            "UX：API 失败提示、账单金额、Dashboard 服务状态准确性",
            "安全：演示密码声明、JWT 密钥生产化、Sonar 质量门",
            "集成测试纳入 CI Gate",
        ])

    # ── 18. Summary ──
    add_title_slide(prs,
        "总结",
        [
            "✅ Go 微服务重构 HIS，管理端门诊闭环可演示",
            "✅ Docker Compose 本地可复现，前后端分离清晰",
            "⚠️ 演示版，部分菜单为占位，支付/安全需生产加固",
            "🔗 github.com/Tangyd893/HIS-Go",
        ],
        bg_color=DARK_BLUE)

    out = r"D:\workspace\coding\HIS-Go\演示\HIS-Go_管理端演示.pptx"
    prs.save(out)
    print(f"✅ 管理端 PPT 已生成: {out}  ({len(prs.slides)} 页)")
    return out


def generate_patient_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── 1. Cover ──
    add_title_slide(prs,
        "HIS-Go 患者端",
        ["微信小程序 + 智能就诊助手", "院外患者服务 · 自助就医 · AI 分诊辅助", "GitHub: Tangyd893/HIS-Go"],
        bg_color=GREEN)

    # ── 2. Positioning ──
    add_content_slide(prs,
        "为什么需要独立患者端？",
        [
            "院内：管理端服务医护人员（挂号、诊疗、发药）",
            "院外：患者端服务就医人群（预约、查报告、慢病、随访）",
            "渠道：H5 浏览器 + 微信小程序（覆盖主流入口）",
            "目标：减少大厅排队、提升就医透明度",
        ])

    # ── 3. Hybrid Architecture ──
    add_content_slide(prs,
        "混合架构：小程序壳 + H5",
        [
            "小程序壳：frontend/mp-webview — 仅负责 <web-view> 容器",
            "业务 H5：frontend/patient — Vue3 + Ant Design Vue + Vite",
            "配置：pages/index/index.js 中 H5_MODE 切换 docker/vite/cloud",
            "Docker 模式：http://127.0.0.1/patient/（Nginx 托管 dist）",
            "一套代码，双端运行：浏览器 + 微信小程序",
        ])

    # ── 4. Tech Stack ──
    add_content_slide(prs,
        "技术选型",
        [
            "前端：Vue 3 · TypeScript · Ant Design Vue · Vue Router · Pinia",
            "路由 base：/patient/（router/index.ts）",
            "布局：顶部 Header + 底部 5 Tab 导航（DefaultLayout.vue）",
            "API：统一走 Gateway http://localhost:8080/api/...",
        ])

    # ── 5. Feature Map ──
    add_content_slide(prs,
        "患者端功能模块",
        [],
        demo_table=(
            ["模块", "路由", "入口"],
            [
                ["首页", "/dashboard", "底部 Tab"],
                ["就诊助手", "/triage", "底部 Tab「助手」"],
                ["预约挂号", "/appointment", "底部 Tab「挂号」"],
                ["我的处方", "/prescription", "底部 Tab「处方」"],
                ["健康档案", "/health-record", "底部 Tab「档案」"],
                ["检查报告", "/report", "首页网格"],
                ["慢病管理", "/chronic", "首页网格"],
                ["我的随访", "/followup", "首页网格"],
            ]
        ),
        note="底部仅 5 Tab，报告/随访需从首页进入——演示时提前熟悉。")

    # ── 6. Backend Services ──
    add_content_slide(prs,
        "demo-patient.yml 服务清单",
        [
            "基础设施：PostgreSQL · Redis（无 RabbitMQ）",
            "核心服务：gateway · auth · user · registration · schedule",
            "患者特色：outpatient（就诊助手）· examination · followup · health_record",
            "静态托管：Nginx 同时提供 /patient/ 与 /api/ 反代",
        ])

    # ── 7. Demo Env ──
    add_content_slide(prs,
        "演示环境与账号",
        [
            "入口：微信开发者工具（推荐）或浏览器 http://localhost/patient",
            "账号：demo-patient / demo123",
            "绑定患者：王小明（patient_001）",
            "⚠️ 微信开发者工具需勾选「不校验合法域名」",
        ],
        note="docker compose -f deploy/compose/demo-patient.yml --env-file deploy/config/demo.env up -d --build")

    # ── 8. RAG Highlight ──
    add_content_slide(prs,
        "RAG + DeepSeek 智能分诊（核心亮点）",
        [
            "知识库：16 类病症（backend/data/triage/knowledge.json）",
            "嵌入：SiliconFlow BGE-M3 语义检索（embeddings.json）",
            "召回：关键词 + 语义双路 Top-K 合并",
            "生成：DeepSeek Chat + 本院科室交集",
            "降级：无 API Key 时关键词模式仍可用",
        ])

    # ── 9. Demo 1: Login ──
    add_content_slide(prs,
        "【现场演示】① 小程序加载与登录",
        [
            "展示微信开发者工具 mp-webview 项目结构",
            "编译运行，显示登录页",
            "demo-patient / demo123 登录",
            "确认进入患者服务中心",
        ],
        demo_table=(
            ["步骤", "环境", "操作"],
            [
                ["1", "微信开发者工具", "展示 mp-webview 项目结构"],
                ["2", "模拟器", "编译运行，显示登录页"],
                ["3", "登录", "demo-patient / demo123"],
                ["4", "首页", "确认进入患者服务中心"],
            ]
        ),
        note="备用方案：浏览器打开 http://localhost/patient")

    # ── 10. Demo 2: Triage ──
    add_content_slide(prs,
        "【现场演示】② AI 就诊助手（重点）",
        [
            "助手 Tab：输入「头痛发热三天，伴有咳嗽」",
            "等待 AI 回复（3-8 秒）",
            "回复内容：推荐科室、就医建议、免责声明",
        ],
        demo_table=(
            ["步骤", "页面", "操作"],
            [
                ["1", "助手 Tab", "输入「头痛发热三天，伴有咳嗽」"],
                ["2", "聊天页", "等待 AI 回复"],
                ["3", "回复内容", "指出推荐科室、就医建议、免责声明"],
            ]
        ),
        note="mode=llm 完整 RAG；mode=keyword 为降级。免责：辅助分诊，不替代医生。")

    # ── 11. Demo 3: Appointment ──
    add_content_slide(prs,
        "【现场演示】③ 预约挂号与处方",
        [
            "挂号 Tab：浏览今日排班（内科/张医生）",
            "处方 Tab：查看处方列表与状态",
            "与管理端共用 schedule/prescription 服务，数据一致",
        ],
        demo_table=(
            ["步骤", "Tab", "操作"],
            [
                ["1", "挂号", "浏览今日排班（内科/张医生）"],
                ["2", "处方", "查看处方列表与状态"],
                ["3", "（可选）", "若管理端刚演示过开方，指新处方"],
            ]
        ))

    # ── 12. Demo 4: Reports ──
    add_content_slide(prs,
        "【现场演示】④ 报告与慢病",
        [
            "首页网格 → 检查报告：展示 seed 报告",
            "首页网格 → 慢病管理：郑十/刘十二签约",
            "档案 Tab：健康档案时间轴",
        ],
        demo_table=(
            ["步骤", "入口", "操作"],
            [
                ["1", "首页网格", "检查报告 → 展示 seed 报告"],
                ["2", "首页网格", "慢病管理 → 郑十/刘十二签约"],
                ["3", "档案 Tab", "健康档案时间轴"],
            ]
        ),
        note="覆盖诊后查结果、慢病长期管理场景。")

    # ── 13. Data Sync ──
    add_content_slide(prs,
        "内外双环 — 与管理端数据联动",
        [
            "管理端开方 → 患者端「我的处方」可见",
            "管理端排班 → 患者端「预约挂号」可见",
            "统一 Gateway :8080 · 统一 JWT 体系",
            "联演示建议：管理端 8 分钟 + 患者端 8 分钟",
        ])

    # ── 14. Real Device ──
    add_content_slide(prs,
        "从模拟器到真机",
        [
            "模拟器：开发者工具本地编译（推荐答辩）",
            "真机：预览扫码，H5 地址改局域网 IP（192.168.x.x/patient/）",
            "上线需配置：微信业务域名 HTTPS、合法 request 域名",
            "当前范围：本地/局域网演示，不上云",
        ])

    # ── 15. Limitations ──
    add_content_slide(prs,
        "诚实口径 — 患者端待完善项",
        [
            "底部 Tab 未覆盖报告/随访（需首页网格进入）",
            "部分路由网关未注册（followup/health-record list）",
            "DEMO_PATIENT_MAP 仅精确映射 demo-patient 账号",
            "就诊助手依赖外部 API Key，需降级预案",
        ])

    # ── 16. Summary ──
    add_title_slide(prs,
        "患者端价值总结",
        [
            "✅ 小程序 + H5 混合架构，一套业务双端复用",
            "✅ 就诊助手 RAG 分诊，差异化亮点",
            "✅ 预约、处方、报告、慢病、随访覆盖院外场景",
            "⚠️ 演示版，上线需域名合规与安全加固",
            "🔗 github.com/Tangyd893/HIS-Go",
        ],
        bg_color=GREEN)

    out = r"D:\workspace\coding\HIS-Go\演示\HIS-Go_患者端演示.pptx"
    prs.save(out)
    print(f"✅ 患者端 PPT 已生成: {out}  ({len(prs.slides)} 页)")
    return out


import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

if __name__ == "__main__":
    generate_admin_ppt()
    generate_patient_ppt()
    print("\n🎉 两个 PPT 全部生成完毕！")
